#!/usr/bin/env python3
"""
Generate a professional PowerPoint presentation for Pothole Detection project.
Run: python create_presentation.py
Output: Pothole_Detection_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
TITLE_COLOR = RGBColor(25, 70, 130)      # Dark blue
ACCENT_COLOR = RGBColor(192, 0, 0)       # Red
LIGHT_BG = RGBColor(240, 248, 255)       # Alice blue
TEXT_COLOR = RGBColor(50, 50, 50)        # Dark gray

def add_title_slide(title, subtitle):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = TITLE_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(32)
    subtitle_p.font.color.rgb = ACCENT_COLOR
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Student info
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
    info_frame = info_box.text_frame
    info_p = info_frame.paragraphs[0]
    info_p.text = "Presented to: Reviewers & Supervisor | Date: February 10, 2026"
    info_p.font.size = Pt(16)
    info_p.font.color.rgb = RGBColor(200, 200, 200)
    info_p.alignment = PP_ALIGN.CENTER

def add_content_slide(title, content_bullets, notes=""):
    """Add a content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = TITLE_COLOR
    
    # Decorative line
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.25), Inches(9), Inches(0))
    shape.line.color.rgb = ACCENT_COLOR
    shape.line.width = Pt(3)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, bullet in enumerate(content_bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_before = Pt(12)
        p.space_after = Pt(12)
    
    # Add notes
    if notes:
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = notes

# Slide 1: Title Slide
add_title_slide("Pothole Detection System", "Automated Road Surface Monitoring using Deep Learning")

# Slide 2: Problem Statement
add_content_slide(
    "Problem Statement",
    [
        "🚗 Road potholes cause vehicle damage, accidents, and injuries",
        "💰 Current manual inspection is expensive and time-consuming",
        "📊 No real-time monitoring or data collection",
        "🛣️ Need for automated, scalable solution"
    ],
    "Emphasize that this is a real-world problem affecting millions. Manual inspection requires significant resources."
)

# Slide 3: Solution Overview
add_content_slide(
    "Solution: Our Approach",
    [
        "✅ Deep Learning-based image classification system",
        "✅ Transfer Learning using pre-trained VGG19 neural network",
        "✅ Binary classification: POTHOLE vs PLAIN ROAD",
        "✅ Optimized for edge deployment (TFLite format)",
        "✅ Real-time inference capability"
    ],
    "Explain that deep learning is perfect for image classification. Transfer learning saves time and resources."
)

# Slide 4: Dataset
add_content_slide(
    "Dataset Overview",
    [
        "📷 Total Images: 712 high-quality road surface photos",
        "🔄 Training Set: 696 images (356 Pothole + 340 Plain)",
        "🧪 Test Set: 16 images (8 Pothole + 8 Plain)",
        "📐 Image Size: Standardized to 224×224 pixels (VGG19 standard)",
        "✨ Data Preprocessing: Normalization, grayscale to RGB conversion, shuffling"
    ],
    "Discuss data collection challenges. Mention balanced classes to avoid bias."
)

# Slide 5: Model Architecture
add_content_slide(
    "Model Architecture: VGG19 Transfer Learning",
    [
        "🏗️ Base Model: VGG19 (pre-trained on ImageNet)",
        "🔒 Frozen Layers: Initial layers frozen to retain learned features",
        "➕ Custom Head: Global Average Pooling → Dense(512) → Dropout(0.5)",
        "                    Dense(256) → Dropout(0.5) → Dense(2, softmax)",
        "⚙️ Loss Function: Categorical Crossentropy",
        "🎯 Optimizer: Adam (learning_rate=0.0005)",
        "🛑 Early Stopping: Prevents overfitting after 10 epochs without improvement"
    ],
    "Explain why VGG19: proven architecture, good for image tasks. Dropout regularization prevents overfitting."
)

# Slide 6: Training Process
add_content_slide(
    "Training & Validation",
    [
        "📊 Batch Size: 32-64 images per batch",
        "🔁 Epochs: Up to 100 (with early stopping)",
        "✔️ Training Metrics: Accuracy, Loss per epoch",
        "🔍 Validation: Monitored after each epoch using hold-out test set",
        "💾 Model Saved: Best weights stored in pothole_model.h5",
        "⚡ Runtime: ~7 seconds per epoch on standard CPU"
    ],
    "Show training progress from logs. Mention that early stopping saves time and improves generalization."
)

# Slide 7: Performance Metrics
add_content_slide(
    "Evaluation Metrics",
    [
        "🎯 Accuracy: Overall correctness of predictions",
        "🔴 Precision (Pothole): True positives / (True + False Positives)",
        "🟢 Recall (Pothole): True positives / (True positives + False Negatives)",
        "📊 F1-Score: Harmonic mean of Precision & Recall (best for imbalanced data)",
        "📈 AUC-ROC: Area under Receiver Operating Characteristic curve",
        "🧮 Matthews Correlation Coefficient (MCC): Robustness metric",
        "📉 Confusion Matrix: Visualizes True/False Positives & Negatives"
    ],
    "Explain why multiple metrics matter. F1 is critical for pothole detection—missing potholes is costly."
)

# Slide 8: Model Outputs
add_content_slide(
    "Key Results",
    [
        "✅ Training Accuracy: ~82% (Epoch 1-20)",
        "✅ Validation Accuracy: ~80%+",
        "✅ F1-Score (Pothole): Strong recall on pothole detection",
        "✅ Low False Negatives: Ensures potholes are not missed",
        "✅ Model Size: Compact H5 format for deployment",
        "✅ Inference Time: <100ms per image"
    ],
    "Share actual numbers from training_output.log. Emphasize low FN rate—missing potholes is unacceptable."
)

# Slide 9: Model Deployment
add_content_slide(
    "Deployment & Edge Computing",
    [
        "📦 Model Formats:",
        "   • pothole_model.h5 (TensorFlow native, full precision)",
        "   • pothole_model.tflite (TensorFlow Lite, optimized for mobile/edge)",
        "🚀 Deployment Targets:",
        "   • Mobile devices (Android/iOS) using TFLite",
        "   • Raspberry Pi / edge computing devices",
        "   • Roadside cameras & vehicle-mounted systems",
        "⚡ Benefits: Real-time inference, low latency, minimal resource usage"
    ],
    "Explain TFLite quantization and model optimization. Show inference speed on edge devices."
)

# Slide 10: Demo
add_content_slide(
    "Live Demo",
    [
        "🖼️ What we'll see: Running inference on sample images",
        "1️⃣ Load a test image from our dataset",
        "2️⃣ Preprocess: Resize to 224×224, normalize",
        "3️⃣ Predict: Model classifies as 'POTHOLE' or 'PLAIN ROAD'",
        "4️⃣ Output: Confidence scores for both classes",
        "5️⃣ Inference Time: Display execution time",
        "",
        "Command: python predictor.py"
    ],
    "Prepare a few test images in advance. Show both pothole and plain road predictions."
)

# Slide 11: Applications & Impact
add_content_slide(
    "Real-World Applications",
    [
        "🛣️ Road Maintenance: Automated pothole detection for city/highway networks",
        "🚗 Vehicle Safety: Early warning system for drivers",
        "💵 Cost Reduction: Reduce manual inspection time by 80%+",
        "📍 Geographic Mapping: Create heat maps of road damage hotspots",
        "🌍 Smart Cities: Integration with IoT and traffic management systems",
        "🔧 Predictive Maintenance: Schedule repairs before accident occurs"
    ],
    "Connect the technology to real-world impact. Mention potential partnerships with municipalities."
)

# Slide 12: Challenges & Solutions
add_content_slide(
    "Challenges & Mitigation",
    [
        "❌ Challenge: Varied lighting conditions in images",
        "   ✅ Solution: Histogram equalization, data augmentation",
        "",
        "❌ Challenge: False positives (shadows/cracks)",
        "   ✅ Solution: Threshold tuning, ensemble methods",
        "",
        "❌ Challenge: Limited test dataset",
        "   ✅ Solution: Continue collecting field data, synthetic augmentation",
        "",
        "❌ Challenge: Real-time processing on edge devices",
        "   ✅ Solution: Model quantization, TFLite optimization, batch processing"
    ],
    "Show awareness of limitations. Discuss mitigation strategies honestly."
)

# Slide 13: Future Enhancements
add_content_slide(
    "Future Work & Improvements",
    [
        "🔍 Extend to multi-class classification (severity: None/Mild/Severe/Critical)",
        "📊 Deploy with mobile app for crowdsourced pothole reporting",
        "🤖 Implement ensemble model (VGG19 + ResNet50 + MobileNet)",
        "💾 Quantize model for ultra-low latency (<50ms)",
        "📡 Integrate with cloud backend for data aggregation & analytics",
        "🗺️ Create interactive web dashboard for visualization",
        "📚 Release as open-source project for broader community use"
    ],
    "Show forward-thinking. Mention scalability and collaboration opportunities."
)

# Slide 14: Code & Resources
add_content_slide(
    "Project Code Structure",
    [
        "📄 main.py: Core training pipeline, model architecture, metrics",
        "📄 predictor.py: Single-image inference script",
        "📄 convert_model.py: TensorFlow to TFLite conversion",
        "📁 My Dataset/: Organized train/test folder structure",
        "📊 training_output.log: Live training metrics and results",
        "💾 pothole_model.h5: Final trained model weights",
        "📱 pothole_model.tflite: Mobile-optimized model"
    ],
    "Briefly explain each file. Offer to share code repository if asked."
)

# Slide 15: Conclusion & Questions
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = TITLE_COLOR

# Main text
conclusion_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4))
conclusion_frame = conclusion_box.text_frame
conclusion_frame.word_wrap = True

p1 = conclusion_frame.paragraphs[0]
p1.text = "Thank You!"
p1.font.size = Pt(54)
p1.font.bold = True
p1.font.color.rgb = RGBColor(255, 255, 255)
p1.alignment = PP_ALIGN.CENTER

p2 = conclusion_frame.add_paragraph()
p2.text = ""
p2.space_before = Pt(24)

p3 = conclusion_frame.add_paragraph()
p3.text = "Pothole Detection System\nAutomated Road Monitoring with Deep Learning"
p3.font.size = Pt(24)
p3.font.color.rgb = ACCENT_COLOR
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(24)

p4 = conclusion_frame.add_paragraph()
p4.text = "\nQuestions?"
p4.font.size = Pt(28)
p4.font.bold = True
p4.font.color.rgb = RGBColor(255, 255, 255)
p4.alignment = PP_ALIGN.CENTER
p4.space_before = Pt(48)

# Save presentation
output_file = "Pothole_Detection_Presentation.pptx"
prs.save(output_file)
print(f"✅ Presentation created: {output_file}")
print(f"📊 Total slides: {len(prs.slides)}")
